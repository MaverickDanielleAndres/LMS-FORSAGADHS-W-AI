from flask import Flask, request, jsonify
from flask_cors import CORS
import os, re, requests, PyPDF2, docx, json
from werkzeug.utils import secure_filename
from pptx import Presentation
import logging

app = Flask(__name__)
CORS(app)

# Configuration
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'doc', 'docx', 'ppt', 'pptx'}
MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_FILE_SIZE

# API Configuration

from dotenv import load_dotenv
load_dotenv()  # Load values from .env file

TOGETHER_API_KEY = os.getenv("TOGETHER_API_KEY")
TOGETHER_API_URL = os.getenv("TOGETHER_API_URL", "https://api.together.xyz/v1/chat/completions")
MODEL_NAME = os.getenv("MODEL_NAME", "meta-llama/Llama-3.3-70B-Instruct-Turbo-Free")

# Logging setup
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_file(file_path, extension):
    """Enhanced text extraction with better error handling"""
    text = ""
    try:
        if extension == 'txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                text = file.read()
                
        elif extension == 'pdf':
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"
                        
        elif extension in ['doc', 'docx']:
            doc = docx.Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
                    
        elif extension in ['ppt', 'pptx']:
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                        
    except Exception as e:
        logger.error(f"Error extracting text from {extension}: {str(e)}")
        text = f"Error reading {extension} file: {str(e)}"
        
    return text.strip()

def create_advanced_prompt(content, settings):
    """Create sophisticated prompts based on exact settings from quiz_maker.php"""
    num_questions = int(settings.get('numQuestions', 10))
    difficulty = settings.get('difficulty', 'mixed')
    question_type = settings.get('questionType', 'mixed')
    points_per_question = float(settings.get('pointsPerQuestion', 2))
    title = settings.get('title', 'Quiz')
    
    # Enhanced difficulty instructions
    difficulty_map = {
        'easy': 'Create straightforward questions testing basic understanding and recall. Focus on simple concepts and direct facts.',
        'medium': 'Create questions requiring analysis and application of concepts. Include some inferential thinking.',
        'hard': 'Create challenging questions requiring critical thinking, synthesis, and deep analysis. Test complex relationships.',
        'mixed': 'Create a balanced distribution: 30% easy (basic recall), 50% medium (analysis/application), 20% hard (synthesis/evaluation)'
    }
    
    # Enhanced question type instructions
    type_map = {
        'mixed': 'Create a variety: 60% multiple choice (4 options each), 25% true/false, 15% short answer questions.',
        'multiple_choice': 'Create ONLY multiple choice questions with exactly 4 options (A, B, C, D). Make distractors plausible but clearly incorrect.',
        'true_false': 'Create ONLY true/false questions. Format as statements with A) True, B) False options.',
        'essay': 'Create essay questions requiring detailed analysis and explanation. Provide clear evaluation criteria.',
        'short_answer': 'Create questions requiring 1-3 sentence responses. Focus on specific facts or brief explanations.'
    }
    
    # Determine content source
    content_instruction = ""
    if content and len(content.strip()) > 50:
        content_instruction = f"Base ALL questions strictly on this provided content:\n\n{content}\n\n"
    else:
        content_instruction = f"Create questions about the topic '{title}'. Generate educational content appropriate for the quiz level.\n\n"

    prompt = f"""You are an expert educational assessment creator. Create EXACTLY {num_questions} high-quality quiz questions.

QUIZ CONFIGURATION:
- Title: {title}
- Number of Questions: {num_questions}
- Difficulty Level: {difficulty}
- Question Type: {question_type}  
- Points per Question: {points_per_question}

{content_instruction}

DIFFICULTY REQUIREMENTS:
{difficulty_map.get(difficulty, difficulty_map['mixed'])}

QUESTION TYPE REQUIREMENTS:
{type_map.get(question_type, type_map['mixed'])}

CRITICAL FORMATTING REQUIREMENTS:
1. Number each question as "QUESTION 1:", "QUESTION 2:", etc.
2. For multiple choice: Use exactly "A) option", "B) option", "C) option", "D) option"
3. For true/false: Use "A) True", "B) False"
4. Add point value after each question: "(Points: {points_per_question})"
5. Include a complete ANSWER KEY section at the end
6. Ensure questions test understanding, not just memorization
7. Make all answer choices plausible for multiple choice questions

QUALITY STANDARDS:
- Each question must be grammatically correct and unambiguous
- Multiple choice distractors must be plausible but clearly incorrect
- Questions should progressively build on concepts
- Avoid trick questions or overly complex wording
- Test different cognitive levels based on difficulty setting

EXAMPLE FORMAT:
QUESTION 1: What is the primary function of photosynthesis in plants? (Points: {points_per_question})
A) To produce oxygen for the atmosphere
B) To convert sunlight into chemical energy
C) To absorb water from the soil
D) To release carbon dioxide

QUESTION 2: The process of photosynthesis occurs primarily in which part of the plant? (Points: {points_per_question})
A) Roots
B) Stems  
C) Leaves
D) Flowers

ANSWER KEY:
1. B ({points_per_question} points) - Photosynthesis converts sunlight into chemical energy (glucose)
2. C ({points_per_question} points) - Photosynthesis occurs mainly in the chloroplasts of leaves

Generate the {num_questions} questions now following this exact format:"""

    return prompt

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({"status": "healthy", "service": "AI Quiz Generator"})

@app.route('/generate-questions', methods=['POST'])
def generate_questions():
    """Enhanced question generation with comprehensive configuration support"""
    try:
        # Extract settings from form data or JSON
        if request.content_type and 'application/json' in request.content_type:
            data = request.get_json()
            settings = {
                'numQuestions': int(data.get('numQuestions', 10)),
                'difficulty': data.get('difficulty', 'mixed'),
                'questionType': data.get('questionType', 'mixed'),
                'pointsPerQuestion': float(data.get('pointsPerQuestion', 2)),
                'title': data.get('title', 'Quiz'),
            }
            content = data.get('content', '')
        else:
            settings = {
                'numQuestions': int(request.form.get('numQuestions', 10)),
                'difficulty': request.form.get('difficulty', 'mixed'),
                'questionType': request.form.get('questionType', 'mixed'),
                'pointsPerQuestion': float(request.form.get('pointsPerQuestion', 2)),
                'title': request.form.get('title', 'Quiz'),
            }
            content = ""
            
            # Handle file upload
            if 'file' in request.files:
                file = request.files['file']
                if file and file.filename and allowed_file(file.filename):
                    filename = secure_filename(file.filename)
                    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                    file.save(file_path)
                    
                    extension = filename.rsplit('.', 1)[1].lower()
                    content = extract_text_from_file(file_path, extension)
                    
                    # Clean up uploaded file
                    os.remove(file_path)
                    
            # Handle material ID (existing materials)
            elif 'material_id' in request.form:
                material_id = request.form.get('material_id')
                content = f"Content from material ID: {material_id}"
                
            # Handle direct prompt/content
            elif 'content' in request.form:
                content = request.form.get('content', '')
        
        # Validate settings
        if settings['numQuestions'] < 1 or settings['numQuestions'] > 50:
            return jsonify({"error": "Number of questions must be between 1 and 50"}), 400
            
        if settings['pointsPerQuestion'] < 0.5 or settings['pointsPerQuestion'] > 10:
            return jsonify({"error": "Points per question must be between 0.5 and 10"}), 400
        
        # Create the enhanced prompt
        prompt = create_advanced_prompt(content, settings)
        
        # Make API call with optimized settings
        headers = {
            "Authorization": f"Bearer {TOGETHER_API_KEY}",
            "Content-Type": "application/json"
        }
        
        # Calculate max tokens based on question count
        estimated_tokens = settings['numQuestions'] * 150 + 500  # ~150 tokens per question + overhead
        max_tokens = min(max(estimated_tokens, 1000), 4000)
        
        payload = {
            "model": MODEL,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.3,  # Lower temperature for more consistent formatting
            "max_tokens": max_tokens,
            "top_p": 0.9,
            "repetition_penalty": 1.1
        }
        
        response = requests.post(TOGETHER_URL, headers=headers, json=payload, timeout=120)
        response.raise_for_status()
        
        ai_response = response.json()['choices'][0]['message']['content']
        
        # Validate response quality
        if not validate_response_quality(ai_response, settings['numQuestions']):
            logger.warning("Response quality validation failed, regenerating...")
            # Try once more with adjusted prompt
            payload['temperature'] = 0.1
            response = requests.post(TOGETHER_URL, headers=headers, json=payload, timeout=120)
            response.raise_for_status()
            ai_response = response.json()['choices'][0]['message']['content']
        
        logger.info(f"Generated {settings['numQuestions']} questions successfully")
        return jsonify({
            "response": ai_response,
            "settings": settings,
            "content_length": len(content) if content else 0
        })
        
    except Exception as e:
        logger.error(f"Error in generate_questions: {str(e)}")
        return jsonify({"error": str(e)}), 500

def validate_response_quality(response, expected_questions):
    """Validate that the response contains the expected number of questions"""
    question_count = len(re.findall(r'QUESTION\s+\d+:', response, re.IGNORECASE))
    answer_key_present = 'ANSWER KEY' in response.upper() or 'Answer Key' in response
    return question_count >= expected_questions * 0.8 and answer_key_present

@app.route('/chat', methods=['POST'])
def chat():
    """Enhanced chat endpoint for AI assistance"""
    try:
        data = request.get_json()
        user_message = data.get('user_message', '')
        
        if not user_message.strip():
            return jsonify({"error": "Empty message"}), 400
            
        # Enhanced system prompt for quiz assistance
        system_prompt = """You are an expert educational AI assistant specializing in quiz and assessment creation for faculty members.

CORE CAPABILITIES:
- Generate questions of any type (multiple choice, true/false, essay, short answer, identification, matching)
- Adapt difficulty levels (easy, medium, hard, mixed)
- Create questions from any educational topic or uploaded material
- Provide educational guidance and best practices
- Suggest improvements to existing questions

RESPONSE GUIDELINES:
- Always provide practical, actionable advice
- When asked to create questions, use proper formatting
- Explain educational reasoning behind suggestions
- Focus on learning objectives and student outcomes
- Be concise but comprehensive
- Offer specific examples when helpful

QUESTION CREATION FORMAT (when requested):
- Use "QUESTION 1:", "QUESTION 2:", etc.
- For multiple choice: A) option, B) option, C) option, D) option
- Include point values: (Points: X)
- Always provide an answer key
- Make questions educationally valuable

EDUCATIONAL PRINCIPLES:
- Test understanding over memorization
- Create plausible but incorrect distractors
- Avoid ambiguous or trick questions
- Align with learning objectives
- Consider cognitive load and difficulty progression

How can I help you create better educational assessments today?"""

        headers = {
            "Authorization": f"Bearer {TOGETHER_API_KEY}",
            "Content-Type": "application/json"
        }
        
        payload = {
            "model": MODEL,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_message}
            ],
            "temperature": 0.7,
            "max_tokens": 2500,
            "top_p": 0.9
        }
        
        response = requests.post(TOGETHER_URL, headers=headers, json=payload, timeout=60)
        response.raise_for_status()
        
        ai_response = response.json()['choices'][0]['message']['content']
        
        return jsonify({"response": ai_response})
        
    except Exception as e:
        logger.error(f"Error in chat: {str(e)}")
        return jsonify({"error": str(e), "response": "I apologize, but I'm experiencing technical difficulties. Please try again."})

@app.route('/get-materials', methods=['GET'])
def get_materials():
    """Get available study materials for quiz generation"""
    try:
        faculty_id = request.args.get('faculty_id')
        subject_code = request.args.get('subject_code')
        
        # This would connect to your database to get materials
        # For now, return a sample response
        materials = [
            {"id": 1, "title": "Chapter 1 - Introduction", "file": "intro.pdf"},
            {"id": 2, "title": "Chapter 2 - Concepts", "file": "concepts.pdf"}
        ]
        
        return jsonify({"materials": materials})
        
    except Exception as e:
        logger.error(f"Error getting materials: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/upload', methods=['POST'])
def upload_and_summarize():
    """File upload and summarization endpoint"""
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400
            
        file = request.files['file']
        if not file or not file.filename:
            return jsonify({"error": "No file selected"}), 400
            
        if not allowed_file(file.filename):
            return jsonify({"error": "File type not supported"}), 400
            
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        try:
            extension = filename.rsplit('.', 1)[1].lower()
            content = extract_text_from_file(file_path, extension)
            
            if not content or len(content.strip()) < 50:
                return jsonify({"error": "Could not extract meaningful content from file"}), 400
            
            # Create summarization prompt
            prompt = f"""Analyze this educational content and provide a comprehensive summary for quiz creation purposes:

{content}

Provide:
1. KEY TOPICS: Main subjects and themes
2. IMPORTANT CONCEPTS: Critical ideas and definitions  
3. FACTUAL CONTENT: Key facts, dates, figures, formulas
4. RELATIONSHIPS: How concepts connect to each other
5. QUIZ POTENTIAL: Areas most suitable for different question types

Focus on content that would be valuable for creating educational assessments."""

            headers = {
                "Authorization": f"Bearer {TOGETHER_API_KEY}",
                "Content-Type": "application/json"
            }
            
            payload = {
                "model": MODEL,
                "messages": [{"role": "user", "content": prompt}],
                "temperature": 0.5,
                "max_tokens": 2000
            }
            
            response = requests.post(TOGETHER_URL, headers=headers, json=payload, timeout=60)
            response.raise_for_status()
            
            ai_response = response.json()['choices'][0]['message']['content']
            
            return jsonify({"response": ai_response, "content_length": len(content)})
            
        finally:
            # Clean up uploaded file
            if os.path.exists(file_path):
                os.remove(file_path)
                
    except Exception as e:
        logger.error(f"Error in upload_and_summarize: {str(e)}")
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    logger.info("Starting Enhanced AI Quiz Generator Server")
    logger.info(f"Upload folder: {UPLOAD_FOLDER}")
    logger.info(f"Allowed extensions: {ALLOWED_EXTENSIONS}")
    app.run(host='0.0.0.0', port=5000, debug=True)