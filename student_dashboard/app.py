from flask import Flask, request, jsonify
from flask_cors import CORS
import os
from werkzeug.utils import secure_filename
import PyPDF2
import docx
import requests
from pptx import Presentation

app = Flask(__name__)
CORS(app)

from dotenv import load_dotenv
load_dotenv()  # Load values from .env file

TOGETHER_API_KEY = os.getenv("TOGETHER_API_KEY")
TOGETHER_API_URL = os.getenv("TOGETHER_API_URL", "https://api.together.xyz/v1/chat/completions")
MODEL_NAME = os.getenv("MODEL_NAME", "meta-llama/Llama-3.3-70B-Instruct-Turbo-Free")

UPLOAD_FOLDER = 'uploads'
# ✅ Only document formats now
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'doc', 'docx', 'ppt', 'pptx'}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB max file size

# Create upload directory if it doesn't exist
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_file_info(file_path, filename):
    """Get basic file information"""
    file_size = os.path.getsize(file_path)
    file_extension = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
    
    return {
        'filename': filename,
        'size': file_size,
        'extension': file_extension,
        'size_mb': round(file_size / (1024 * 1024), 2)
    }

def extract_text_content(file_path, file_extension):
    """Extract text content from supported document formats"""
    extracted_text = ""
    
    try:
        if file_extension == 'txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                extracted_text = f.read()
                
        elif file_extension == 'pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        extracted_text += text + "\n"
                        
        elif file_extension == 'docx':
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                extracted_text += para.text + '\n'
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        extracted_text += cell.text + '\t'
                    extracted_text += '\n'
                    
        elif file_extension in ['pptx']:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                extracted_text += f"\n--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        extracted_text += shape.text + '\n'
                    # Extract text from tables in slides
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                extracted_text += cell.text + '\t'
                            extracted_text += '\n'
                            
        elif file_extension in ['ppt', 'doc']:
            # For older formats, return a message suggesting conversion
            extracted_text = f"This is a {file_extension.upper()} file. For best results, please convert to {file_extension.upper()}X format."
                    
    except Exception as e:
        extracted_text = f"Error reading {file_extension.upper()} content: {str(e)}"

    return extracted_text

def generate_summary_prompt(file_info, content=None):
    """Generate appropriate summary prompt based on file type"""
    filename = file_info['filename']
    extension = file_info['extension']
    size_mb = file_info['size_mb']
    
    if content and content.strip():
        prompt = f"""Please provide a comprehensive summary and analysis of this {extension.upper()} file:

File: {filename} ({size_mb} MB)

Content:
{content}

Please provide:
1. Executive Summary
2. Key Points and Main Topics
3. Structure and Organization
4. Important Details or Data
5. Conclusions or Recommendations (if applicable)
"""
    else:
        prompt = f"""I have a {extension.upper()} file named "{filename}" ({size_mb} MB) that couldn't be read properly. 
Please provide general information about:
1. What {extension.upper()} files typically contain
2. Common use cases for this file format
3. Suggestions for accessing the content
4. Alternative methods to view or convert the file
"""
    
    return prompt

def generate_questions_prompt(file_info, content=None, num_questions=10):
    """Enhanced question prompt matching the exact format requested"""
    filename = file_info['filename']
    extension = file_info['extension']
    size_mb = file_info['size_mb']
    
    if content and content.strip():
        prompt = f"""Based on the content of this {extension.upper()} file, create {num_questions} comprehensive review questions for students:

File: {filename} ({size_mb} MB)

Content:
{content}

Please create exactly {num_questions} multiple-choice questions that:
1. Cover the main topics and key concepts from the document
2. Test different levels of understanding (basic recall, comprehension, analysis, application)
3. Are clearly written and unambiguous
4. Have 4 answer choices each (A, B, C, D)
5. Include a variety of question types (factual, conceptual, analytical)

Format your response EXACTLY like this example:

QUESTION 1: What is one of the primary benefits of integrating an AI-generated reviewer in the Learning Management System (LMS) for Rizal High School?

A) Reducing the need for internet connectivity
B) Enhancing student concentration and retention
C) Increasing the use of paper materials
D) Reducing the need for teachers

QUESTION 2: [Your question text here]

A) [Option A]
B) [Option B]
C) [Option C]
D) [Option D]

Continue this exact format for all {num_questions} questions.

After all questions, provide:

ANSWER KEY:
1. [Correct letter]
2. [Correct letter]
3. [Correct letter]
4. [Correct letter]
5. [Correct letter]
6. [Correct letter]
7. [Correct letter]
8. [Correct letter]
9. [Correct letter]
10. [Correct letter]

IMPORTANT FORMATTING RULES:
- Use "QUESTION X:" (with colon) at the start of each question
- Leave one blank line after each question text
- Format options as "A) Option text" (with closing parenthesis and space)
- Leave one blank line between each question block
- Make sure questions test important concepts from the material
- Ensure all questions are educational and academically appropriate"""
    else:
        prompt = f"""I have a {extension.upper()} file named "{filename}" ({size_mb} MB) that couldn't be read properly. 
Please create {num_questions} general review questions about typical {extension.upper()} document content and file format knowledge using the exact format specified above."""
    
    return prompt

def chat_with_ai(user_message):
    headers = {
        "Authorization": f"Bearer {TOGETHER_API_KEY}",
        "Content-Type": "application/json"
    }

    data = {
        "model": MODEL_NAME,
        "messages": [{"role": "user", "content": user_message}],
        "temperature": 0.7,
        "max_tokens": 2048
    }

    try:
        response = requests.post(TOGETHER_API_URL, headers=headers, json=data)
        response.raise_for_status()

        result = response.json()

        # Try multiple ways to extract content
        ai_response = ''
        if 'choices' in result and len(result['choices']) > 0:
            ai_response = result['choices'][0].get('message', {}).get('content', '')
            if not ai_response:
                ai_response = result['choices'][0].get('text', '')

        if not ai_response:
            return jsonify({"error": "Failed to parse AI response."}), 500

        return jsonify({"response": ai_response})

    except requests.exceptions.RequestException as e:
        return jsonify({"error": f"API request failed: {str(e)}", "details": response.text}), 500
    except KeyError as e:
        return jsonify({"error": f"Unexpected API response format: {str(e)}"}), 500
    except Exception as e:
        return jsonify({"error": f"Unexpected error: {str(e)}"}), 500


@app.route('/chat', methods=['POST'])
def chat():
    try:
        user_message = request.json.get('user_message')
        if not user_message:
            return jsonify({"error": "No message provided"}), 400

        return chat_with_ai(user_message)
    except Exception as e:
        return jsonify({"error": f"Chat error: {str(e)}"}), 500

@app.route('/upload', methods=['POST'])
def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No file part"}), 400
        file = request.files['file']

        if file.filename == '':
            return jsonify({"error": "No selected file"}), 400

        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)

            # Get file information
            file_info = get_file_info(file_path, filename)
            file_extension = file_info['extension']
            
            # Extract content from document
            content = extract_text_content(file_path, file_extension)
            
            # Generate appropriate summary prompt
            summary_prompt = generate_summary_prompt(file_info, content)
            
            # Clean up the uploaded file
            try:
                os.remove(file_path)
            except:
                pass
            
            # Get AI summary
            return chat_with_ai(summary_prompt)

        return jsonify({"error": "Invalid file format. Supported formats: TXT, PDF, DOC, DOCX, PPT, PPTX"}), 400
    except Exception as e:
        return jsonify({"error": f"Upload error: {str(e)}"}), 500

@app.route('/analyze', methods=['POST'])
def analyze_file():
    """Enhanced analysis endpoint with custom prompts"""
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No file part"}), 400
        
        file = request.files['file']
        analysis_type = request.form.get('analysis_type', 'summary')
        custom_prompt = request.form.get('custom_prompt', '')

        if file.filename == '':
            return jsonify({"error": "No selected file"}), 400

        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)

            # Get file information
            file_info = get_file_info(file_path, filename)
            file_extension = file_info['extension']
            
            # Extract content from document
            content = extract_text_content(file_path, file_extension)
            
            # Generate custom analysis prompt
            if custom_prompt:
                if content:
                    analysis_prompt = f"{custom_prompt}\n\nFile: {filename}\nContent:\n{content}"
                else:
                    analysis_prompt = f"{custom_prompt}\n\nFile: {filename} ({file_extension.upper()}, {file_info['size_mb']} MB)"
            else:
                # Use predefined analysis types
                if analysis_type == 'detailed':
                    analysis_prompt = f"Provide a detailed, comprehensive analysis of this file including technical aspects, structure, and in-depth content breakdown:\n\nFile: {filename}\n"
                    if content:
                        analysis_prompt += f"Content:\n{content}"
                elif analysis_type == 'technical':
                    analysis_prompt = f"Provide a technical analysis focusing on format, structure, metadata, and technical specifications:\n\nFile: {filename} ({file_extension.upper()}, {file_info['size_mb']} MB)\n"
                    if content:
                        analysis_prompt += f"Content:\n{content}"
                elif analysis_type == 'creative':
                    analysis_prompt = f"Provide a creative interpretation and analysis of this file, including potential creative uses and artistic perspectives:\n\nFile: {filename}\n"
                    if content:
                        analysis_prompt += f"Content:\n{content}"
                else:  # default summary
                    analysis_prompt = generate_summary_prompt(file_info, content)
            
            # Clean up the uploaded file
            try:
                os.remove(file_path)
            except:
                pass
            
            # Get AI analysis
            return chat_with_ai(analysis_prompt)

        return jsonify({"error": "Invalid file format"}), 400
    except Exception as e:
        return jsonify({"error": f"Analysis error: {str(e)}"}), 500

@app.route('/generate-questions', methods=['POST'])
def generate_questions():
    """Generate review questions based on uploaded document with custom options"""
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No file part"}), 400
        
        file = request.files['file']
        
        # Get optional parameters
        num_questions = int(request.form.get('num_questions', 10))
        difficulty = request.form.get('difficulty', 'mixed')
        question_type = request.form.get('question_type', 'mixed')

        if file.filename == '':
            return jsonify({"error": "No selected file"}), 400

        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)

            # Get file information
            file_info = get_file_info(file_path, filename)
            file_extension = file_info['extension']
            
            # Extract content from document
            content = extract_text_content(file_path, file_extension)
            
            # Generate questions prompt with custom options
            questions_prompt = generate_questions_prompt(file_info, content, num_questions)
            
            # Add difficulty and type specifications to the prompt
            if difficulty != 'mixed':
                questions_prompt += f"\n\nIMPORTANT: Make all questions {difficulty.upper()} difficulty level."
            
            if question_type != 'mixed':
                questions_prompt += f"\n\nIMPORTANT: Focus primarily on {question_type.upper()} questions."
            
            # Clean up the uploaded file
            try:
                os.remove(file_path)
            except:
                pass
            
            # Get AI generated questions
            return chat_with_ai(questions_prompt)

        return jsonify({"error": "Invalid file format. Supported formats: TXT, PDF, DOC, DOCX, PPT, PPTX"}), 400
    except Exception as e:
        return jsonify({"error": f"Question generation error: {str(e)}"}), 500

@app.route('/generate-custom-questionnaire', methods=['POST'])
def generate_custom_questionnaire():
    """Generate questionnaire with completely custom parameters"""
    try:
        data = request.get_json()
        
        if not data:
            return jsonify({"error": "No data provided"}), 400
        
        # Extract parameters
        topic = data.get('topic', 'General Knowledge')
        num_questions = data.get('num_questions', 10)
        difficulty = data.get('difficulty', 'mixed')
        subject_area = data.get('subject_area', 'General')
        specific_topics = data.get('specific_topics', [])
        
        # Build custom prompt
        prompt = f"""Create {num_questions} multiple-choice questions about {topic} for {subject_area}.

REQUIREMENTS:
- Use the exact format shown in this example:

QUESTION 1: What is one of the primary benefits of integrating an AI-generated reviewer in the Learning Management System (LMS) for Rizal High School?

A) Reducing the need for internet connectivity
B) Enhancing student concentration and retention
C) Increasing the use of paper materials
D) Reducing the need for teachers

- Continue with QUESTION 2, QUESTION 3, etc.
- Each question should have exactly 4 options (A, B, C, D)
- Leave blank lines between questions for readability
- Questions should be {difficulty} difficulty level
"""
        
        if specific_topics:
            prompt += f"\n- Focus on these specific topics: {', '.join(specific_topics)}"
        
        prompt += f"""

After all questions, provide:

ANSWER KEY:
1. [Correct letter]
2. [Correct letter]
...
{num_questions}. [Correct letter]

Make sure all questions are educationally valuable and test important concepts."""
        
        return chat_with_ai(prompt)
    except Exception as e:
        return jsonify({"error": f"Custom questionnaire error: {str(e)}"}), 500

@app.route('/supported-formats', methods=['GET'])
def supported_formats():
    return jsonify({
        "documents": list(ALLOWED_EXTENSIONS),
        "all_supported": list(ALLOWED_EXTENSIONS),
        "max_file_size_mb": 100
    })

@app.route('/health', methods=['GET'])
def health_check():
    return jsonify({
        "status": "healthy",
        "service": "Document Analysis API",
        "supported_formats": {
            "documents": list(ALLOWED_EXTENSIONS)
        },
        "max_file_size_mb": 100,
        "features": [
            "Document summarization",
            "Document analysis", 
            "Question generation",
            "Custom questionnaire generation",
            "Custom analysis prompts",
            "Document content extraction",
            "Multiple difficulty levels",
            "Flexible question formatting"
        ]
    })
if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)